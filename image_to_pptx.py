from pptx.util import Inches
from wand.compat import text
from wand.image import Image
from pptx import Presentation
from os import chdir, listdir
import pytest;

# Counts number of images to be itirated over
img_count = len(listdir('./input_img'));

# pptx init
prs = Presentation();
bullet_slide_layout = prs.slide_layouts[1]

# Resizes image to smaller dimensions adds adds logo
def img_logo(i):
    with Image(filename = './input_img/image{}.jpg'.format(i+1)) as img1:
        with Image(filename = 'nike_black.png') as img2:
            img1.transform(resize = '5.5%');
            img2.transform(resize = '2%');
            img1.transparentize(0.35)
            #img2.transparentize(0.35)
            img1.composite(image=img2, left=0, top=0)
            img1.save(filename = './output_img/image-edited{}.jpg'.format(i+1));

# Iterates over each slide and adds logo and heading to it
def create_side(prs, bullet_slide_layout, i):
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    body_shape = shapes.placeholders[1]
    # Adding Image to each slide
    left = Inches(1);
    top = Inches(2.5);
    slide.shapes.add_picture('./output_img/image-edited{}.jpg'.format(i+1), left, top);
    tf = body_shape.text_frame
    # Add Heading
    title_shape = shapes.title
    title_shape.text = "Sample Heading {}".format(i+1);
    # Add Subheading
    tf.text = "Sample Subheading {}".format(i+1);
for i in range(img_count):
    img_logo(i);
    create_side(prs, bullet_slide_layout, i)
    prs.save('as.pptx');